import os
import uuid
from pathlib import Path

from dotenv import load_dotenv
from PyPDF2 import PdfReader
from pptx import Presentation
import requests
from bs4 import BeautifulSoup
import chromadb
from openai import OpenAI
from backend.storage import save_raw_file, save_extracted_text, notebook_dir

load_dotenv()
client = OpenAI()  # reads OPENAI_API_KEY from environment automatically
DATA_DIR = Path("data/users")


# Directory helper

def get_chroma_path(username: str, nb_id: str) -> str:
    path = notebook_dir(username, nb_id) / "chroma"
    path.mkdir(parents=True, exist_ok=True)
    return str(path)


# Chunking

def chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> list:
    """Split text into overlapping word-based chunks."""
    words = text.split()
    chunks = []
    for i in range(0, len(words), chunk_size - overlap):
        chunk = " ".join(words[i : i + chunk_size])
        if chunk:
            chunks.append(chunk)
    return chunks


# Embedding + ChromaDB storage

def embed_and_store(chunks: list, source_name: str, nb_id: str, username: str):
    """Embed chunks using OpenAI and store in ChromaDB."""
    chroma_path = get_chroma_path(username, nb_id)
    chroma_client = chromadb.PersistentClient(path=chroma_path)
    collection = chroma_client.get_or_create_collection(
        name="sources",
        metadata={"hnsw:space": "cosine"},
    )

    # Embed in batches of 100 to stay within API limits
    batch_size = 100
    for batch_start in range(0, len(chunks), batch_size):
        batch = chunks[batch_start : batch_start + batch_size]

        response = client.embeddings.create(
            model="text-embedding-3-small",
            input=batch,
        )
        embeddings = [item.embedding for item in response.data]

        collection.add(
            documents=batch,
            embeddings=embeddings,
            metadatas=[
                {"source": source_name, "chunk": batch_start + i}
                for i in range(len(batch))
            ],
            ids=[str(uuid.uuid4()) for _ in batch],
        )


# Text extraction

def extract_pdf(filepath: str) -> str:
    """Extract text from a PDF file."""
    reader = PdfReader(filepath)
    pages = []
    for page_num, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        if text.strip():
            pages.append(f"[Page {page_num + 1}]\n{text}")
    return "\n\n".join(pages)


def extract_pptx(filepath: str) -> str:
    """Extract text from a PPTX file."""
    prs = Presentation(filepath)
    texts = []
    for slide_num, slide in enumerate(prs.slides):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
        if slide_texts:
            texts.append(f"[Slide {slide_num + 1}]\n" + "\n".join(slide_texts))
    return "\n\n".join(texts)


def extract_txt(filepath: str) -> str:
    """Read a plain text file."""
    with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def fetch_url(url: str) -> str:
    """Fetch and extract readable text from a web page."""
    headers = {"User-Agent": "Mozilla/5.0"}
    response = requests.get(url, timeout=15, headers=headers)
    response.raise_for_status()
    soup = BeautifulSoup(response.text, "html.parser")

    # Remove script and style tags
    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.decompose()

    return soup.get_text(separator=" ", strip=True)


# Public ingest functions

def ingest_file(filepath: str, nb_id: str, username: str) -> str:
    """Extract, chunk, embed, and store a file into the notebook vector store."""
    ext = Path(filepath).suffix.lower()
    filename = Path(filepath).name

    try:
        if ext == ".pdf":
            text = extract_pdf(filepath)
        elif ext == ".pptx":
            text = extract_pptx(filepath)
        elif ext == ".txt":
            text = extract_txt(filepath)
        else:
            return f"Unsupported file type: {ext}. Supported: .pdf, .pptx, .txt"

        if not text.strip():
            return f"Warning: no text could be extracted from {filename}"

        # Save raw file and extracted text
        with open(filepath, "rb") as f:
            save_raw_file(username, nb_id, filename, f.read())
        save_extracted_text(username, nb_id, filename, text)

        # Chunk and embed
        chunks = chunk_text(text)
        embed_and_store(chunks, filename, nb_id, username)

        return f"Ingested {filename} ({len(chunks)} chunks)"

    except Exception as e:
        return f"Error ingesting {filename}: {str(e)}"


def ingest_url(url: str, nb_id: str, username: str) -> str:
    """Fetch, chunk, embed, and store a web page into the notebook vector store."""
    try:
        text = fetch_url(url)

        if not text.strip():
            return f"Warning: no text found at {url}"

        save_extracted_text(username, nb_id, url.replace("/", "_"), text)

        chunks = chunk_text(text)
        embed_and_store(chunks, url, nb_id, username)

        return f"Ingested URL: {url} ({len(chunks)} chunks)"

    except requests.exceptions.RequestException as e:
        return f"Error fetching URL {url}: {str(e)}"
    except Exception as e:
        return f"Error ingesting URL {url}: {str(e)}"