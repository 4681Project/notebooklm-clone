"""
ingestion.py
------------
Ingestion pipeline: extract text from PDF / PPTX / TXT / URL,
chunk the text, embed it, and store it in a per-notebook ChromaDB collection.

Chunking strategies available:
  - "fixed"      : fixed token-count windows with overlap
  - "sentence"   : sentence-boundary-aware sliding window
  - "recursive"  : recursive character splitter (LangChain-style)

Usage:
    result = ingest_source(username, notebook_id, source_path_or_url)
"""

import hashlib
import re
import os
from pathlib import Path
from typing import Literal, Optional
from urllib.parse import urlparse

import chromadb
from chromadb.utils.embedding_functions import OpenAIEmbeddingFunction

from backend import storage

# ---------------------------------------------------------------------------
# Config
# ---------------------------------------------------------------------------

OPENAI_API_KEY = os.getenv("OPENAI_API_KEY", "")
CHUNK_SIZE = int(os.getenv("CHUNK_SIZE", "800"))       # tokens / chars
CHUNK_OVERLAP = int(os.getenv("CHUNK_OVERLAP", "150"))
EMBED_MODEL = os.getenv("EMBED_MODEL", "text-embedding-3-small")
ChunkStrategy = Literal["fixed", "sentence", "recursive"]
DEFAULT_CHUNK_STRATEGY: ChunkStrategy = "recursive"


# ---------------------------------------------------------------------------
# Embedding function
# ---------------------------------------------------------------------------

def _get_embedding_fn():
    return OpenAIEmbeddingFunction(
        api_key=OPENAI_API_KEY,
        model_name=EMBED_MODEL,
    )


# ---------------------------------------------------------------------------
# ChromaDB helpers
# ---------------------------------------------------------------------------

def _get_collection(username: str, notebook_id: str):
    chroma_path = storage.get_chroma_path(username, notebook_id)
    client = chromadb.PersistentClient(path=chroma_path)
    collection = client.get_or_create_collection(
        name="sources",
        embedding_function=_get_embedding_fn(),
        metadata={"hnsw:space": "cosine"},
    )
    return collection


# ---------------------------------------------------------------------------
# Text extraction
# ---------------------------------------------------------------------------

def extract_text_from_pdf(file_path: str | Path) -> str:
    """Extract text from a PDF using pdfplumber."""
    import pdfplumber
    texts = []
    with pdfplumber.open(str(file_path)) as pdf:
        for page in pdf.pages:
            t = page.extract_text()
            if t:
                texts.append(t)
    return "\n\n".join(texts)


def extract_text_from_pptx(file_path: str | Path) -> str:
    """Extract text from a PowerPoint file."""
    from pptx import Presentation
    prs = Presentation(str(file_path))
    texts = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_texts = [f"[Slide {slide_num}]"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        slide_texts.append(line)
        texts.append("\n".join(slide_texts))
    return "\n\n".join(texts)


def extract_text_from_url(url: str) -> str:
    """Extract readable text from a web URL using trafilatura."""
    import trafilatura
    downloaded = trafilatura.fetch_url(url)
    if downloaded is None:
        raise ValueError(f"Could not fetch URL: {url}")
    text = trafilatura.extract(downloaded, include_links=False, include_tables=True)
    if not text:
        raise ValueError(f"No readable text found at URL: {url}")
    return text


def extract_text(source: str | Path) -> tuple[str, str]:
    """
    Dispatch extraction based on source type.
    Returns (text, source_name).
    """
    src = str(source)
    parsed = urlparse(src)
    if parsed.scheme in ("http", "https"):
        text = extract_text_from_url(src)
        name = parsed.netloc + parsed.path.replace("/", "_")
        return text, name

    p = Path(src)
    suffix = p.suffix.lower()
    if suffix == ".pdf":
        return extract_text_from_pdf(p), p.name
    elif suffix in (".pptx", ".ppt"):
        return extract_text_from_pptx(p), p.name
    elif suffix in (".txt", ".md"):
        return p.read_text(encoding="utf-8", errors="ignore"), p.name
    else:
        raise ValueError(f"Unsupported file type: {suffix}")


# ---------------------------------------------------------------------------
# Chunking strategies
# ---------------------------------------------------------------------------

def _fixed_chunks(text: str, size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    """Split on character count with overlap."""
    chunks, start = [], 0
    while start < len(text):
        end = start + size
        chunks.append(text[start:end])
        start += size - overlap
    return [c.strip() for c in chunks if c.strip()]


def _sentence_chunks(text: str, size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    """Split on sentence boundaries, then group into windows."""
    sentences = re.split(r'(?<=[.!?])\s+', text)
    chunks, current, current_len = [], [], 0
    for sent in sentences:
        slen = len(sent)
        if current_len + slen > size and current:
            chunks.append(" ".join(current))
            # keep overlap sentences
            overlap_sents = []
            ol = 0
            for s in reversed(current):
                ol += len(s)
                overlap_sents.insert(0, s)
                if ol >= overlap:
                    break
            current = overlap_sents
            current_len = sum(len(s) for s in current)
        current.append(sent)
        current_len += slen
    if current:
        chunks.append(" ".join(current))
    return [c.strip() for c in chunks if c.strip()]


def _recursive_chunks(text: str, size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    """
    Recursive character splitter — tries paragraph, then sentence,
    then word boundaries before hard-splitting.
    """
    separators = ["\n\n", "\n", ". ", " ", ""]

    def _split(t: str, seps: list[str]) -> list[str]:
        if len(t) <= size:
            return [t] if t.strip() else []
        if not seps:
            return _fixed_chunks(t, size, overlap)

        sep = seps[0]
        parts = t.split(sep)
        chunks, current = [], ""
        for part in parts:
            candidate = current + (sep if current else "") + part
            if len(candidate) <= size:
                current = candidate
            else:
                if current:
                    chunks.extend(_split(current, seps[1:]))
                current = part
        if current:
            chunks.extend(_split(current, seps[1:]))
        return chunks

    raw = _split(text, separators)
    # Add overlap between adjacent chunks
    result = []
    for i, chunk in enumerate(raw):
        if i > 0 and overlap > 0:
            prev_tail = raw[i - 1][-overlap:]
            chunk = prev_tail + " " + chunk
        result.append(chunk.strip())
    return [c for c in result if c]


def chunk_text(text: str, strategy: ChunkStrategy = DEFAULT_CHUNK_STRATEGY) -> list[str]:
    if strategy == "fixed":
        return _fixed_chunks(text)
    elif strategy == "sentence":
        return _sentence_chunks(text)
    else:
        return _recursive_chunks(text)


# ---------------------------------------------------------------------------
# Main ingest function
# ---------------------------------------------------------------------------

def ingest_source(
    username: str,
    notebook_id: str,
    source: str | Path,
    chunk_strategy: ChunkStrategy = DEFAULT_CHUNK_STRATEGY,
    raw_bytes: Optional[bytes] = None,
) -> dict:
    """
    Full ingestion pipeline for one source.
    Returns a summary dict with stats.
    """
    # 1. Extract text
    text, source_name = extract_text(source)

    # 2. Persist raw file (if bytes provided, i.e. an upload)
    if raw_bytes is not None:
        storage.save_raw_file(username, notebook_id, source_name, raw_bytes)
    
    # 3. Persist extracted text
    storage.save_extracted_text(username, notebook_id, source_name, text)

    # 4. Chunk
    chunks = chunk_text(text, strategy=chunk_strategy)

    # 5. Build IDs and metadata
    ids, documents, metadatas = [], [], []
    for i, chunk in enumerate(chunks):
        chunk_id = hashlib.md5(f"{source_name}_{i}_{chunk[:40]}".encode()).hexdigest()
        ids.append(chunk_id)
        documents.append(chunk)
        metadatas.append({
            "source": source_name,
            "chunk_index": i,
            "total_chunks": len(chunks),
            "strategy": chunk_strategy,
        })

    # 6. Upsert into ChromaDB
    collection = _get_collection(username, notebook_id)
    # Batch upsert in groups of 100 to avoid request-size limits
    batch_size = 100
    for start in range(0, len(ids), batch_size):
        end = start + batch_size
        collection.upsert(
            ids=ids[start:end],
            documents=documents[start:end],
            metadatas=metadatas[start:end],
        )

    return {
        "source_name": source_name,
        "char_count": len(text),
        "chunk_count": len(chunks),
        "strategy": chunk_strategy,
    }


def delete_source(username: str, notebook_id: str, source_name: str) -> int:
    """Remove all chunks for a source from the vector store. Returns deleted count."""
    collection = _get_collection(username, notebook_id)
    results = collection.get(where={"source": source_name})
    if results and results["ids"]:
        collection.delete(ids=results["ids"])
        return len(results["ids"])
    return 0


def list_indexed_sources(username: str, notebook_id: str) -> list[str]:
    """Return unique source names currently indexed in ChromaDB."""
    collection = _get_collection(username, notebook_id)
    results = collection.get(include=["metadatas"])
    sources = {m["source"] for m in results["metadatas"]} if results["metadatas"] else set()
    return sorted(sources)